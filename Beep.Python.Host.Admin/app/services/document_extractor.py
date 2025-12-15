"""
Document Extraction Service

Standalone service for extracting text from various document formats:
- PDF (PyPDF2, pypdf, PyMuPDF)
- DOCX (python-docx)
- XLSX (openpyxl)
- TXT, MD, CSV, JSON, HTML, XML, etc.

This service can be used independently or integrated with RAG/LLM systems.

Uses a dedicated virtual environment for document extraction libraries.
"""
import os
import io
import subprocess
import sys
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    """Result of document extraction"""
    success: bool
    text: str
    metadata: Dict[str, Any]
    error: Optional[str] = None


class DocumentExtractor:
    """
    Standalone document extraction service.
    
    Supports:
    - PDF: .pdf
    - Word: .docx
    - Excel: .xlsx, .xls
    - Text: .txt, .md, .csv, .json, .html, .xml, .yaml, .yml
    - Code: .py, .js, .ts, .css, etc.
    """
    
    # Supported file extensions
    SUPPORTED_EXTENSIONS = {
        # Documents
        'pdf': 'PDF',
        'docx': 'Word',
        'doc': 'Word (legacy)',
        'xlsx': 'Excel',
        'xls': 'Excel (legacy)',
        'pptx': 'PowerPoint',
        'ppt': 'PowerPoint (legacy)',
        # Images (OCR)
        'png': 'PNG Image',
        'jpg': 'JPEG Image',
        'jpeg': 'JPEG Image',
        'gif': 'GIF Image',
        'bmp': 'BMP Image',
        'tiff': 'TIFF Image',
        'tif': 'TIFF Image',
        'webp': 'WebP Image',
        # Text files
        'txt': 'Text',
        'md': 'Markdown',
        'csv': 'CSV',
        'json': 'JSON',
        'html': 'HTML',
        'xml': 'XML',
        'yaml': 'YAML',
        'yml': 'YAML',
        # Code files
        'py': 'Python',
        'js': 'JavaScript',
        'ts': 'TypeScript',
        'css': 'CSS',
        'java': 'Java',
        'cpp': 'C++',
        'c': 'C',
        'h': 'C Header',
        'rs': 'Rust',
        'go': 'Go',
    }
    
    def __init__(self, use_dedicated_env: bool = True):
        """
        Initialize document extractor.
        
        Args:
            use_dedicated_env: If True, use dedicated environment for extraction libraries
        """
        self.use_dedicated_env = use_dedicated_env
        
        if use_dedicated_env:
            from app.services.document_extraction_environment import get_doc_extraction_env
            self.env_mgr = get_doc_extraction_env()
        else:
            self.env_mgr = None
        
        self._libraries_available = self._check_libraries()
        
        # OCR engine tracker
        from app.services.ocr_engine_tracker import get_ocr_engine_tracker, OCREngineType
        self._ocr_tracker = get_ocr_engine_tracker()
        # Default to EasyOCR if no engine is selected
        if self._ocr_tracker.get_active_engine() is None:
            self._ocr_tracker.activate_engine(OCREngineType.EASYOCR)
    
    def _get_python_executable(self) -> Optional[Path]:
        """Get Python executable from virtual environment using environment manager"""
        if not self.env_mgr:
            return None
        
        try:
            python_path = self.env_mgr._get_python_path()
            if python_path and python_path.exists():
                logger.debug(f"Using Python from document extraction virtual environment: {python_path}")
                return python_path
            else:
                logger.warning("Virtual environment Python not found, using system Python")
                return None
        except Exception as e:
            logger.error(f"Error getting Python executable: {e}")
            return None
    
    def refresh_libraries(self):
        """Refresh library availability checks (useful after installing packages)"""
        self._libraries_available = self._check_libraries()
    
    def _check_libraries(self) -> Dict[str, bool]:
        """Check which extraction libraries are available"""
        libraries = {
            'pypdf': False,
            'pymupdf': False,
            'python_docx': False,
            'openpyxl': False,
            'xlrd': False,
            'easyocr': False,
            'pillow': False,
        }
        
        # Determine which Python to use - get dynamically from environment manager
        python_exe_path = self._get_python_executable()
        python_exe = str(python_exe_path) if python_exe_path else sys.executable
        
        # Check libraries using the appropriate Python
        def check_import(module_name, test_code, timeout=5):
            try:
                result = subprocess.run(
                    [python_exe, '-c', test_code],
                    capture_output=True,
                    text=True,
                    timeout=timeout
                )
                if result.returncode != 0:
                    # Log error for debugging
                    logger.debug(f"Failed to import {module_name} in {python_exe}: {result.stderr[:200]}")
                return result.returncode == 0
            except subprocess.TimeoutExpired:
                logger.warning(f"Timeout checking {module_name} import (timeout={timeout}s)")
                return False
            except Exception as e:
                logger.debug(f"Exception checking {module_name}: {e}")
                return False
        
        # Check PyPDF2/pypdf
        libraries['pypdf'] = check_import('pypdf', 'import pypdf') or check_import('PyPDF2', 'import PyPDF2')
        
        # Check PyMuPDF
        libraries['pymupdf'] = check_import('fitz', 'import fitz')
        
        # Check python-docx
        libraries['python_docx'] = check_import('docx', 'import docx')
        
        # Check openpyxl
        libraries['openpyxl'] = check_import('openpyxl', 'import openpyxl')
        
        # Check xlrd
        libraries['xlrd'] = check_import('xlrd', 'import xlrd')
        
        # Check Pillow (PIL)
        libraries['pillow'] = check_import('PIL', 'from PIL import Image')
        
        # Check easyocr - it can take time to import on first run (loads models)
        libraries['easyocr'] = check_import('easyocr', 'import easyocr', timeout=30)
        
        # Check tesseract
        libraries['tesseract'] = check_import('pytesseract', 'import pytesseract', timeout=10)
        
        # Check paddleocr
        libraries['paddleocr'] = check_import('paddleocr', 'from paddleocr import PaddleOCR', timeout=30)
        
        return libraries
    
    def is_supported(self, filename: str) -> bool:
        """Check if file type is supported"""
        ext = self._get_extension(filename)
        return ext in self.SUPPORTED_EXTENSIONS
    
    def _get_extension(self, filename: str) -> str:
        """Get file extension in lowercase"""
        return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    def extract(self, file_path: Optional[str] = None, 
                file_content: Optional[bytes] = None,
                filename: Optional[str] = None) -> ExtractionResult:
        """
        Extract text from a document file.
        
        Args:
            file_path: Path to the file (if reading from disk)
            file_content: File content as bytes (if reading from memory)
            filename: Filename for determining file type
        
        Returns:
            ExtractionResult with extracted text and metadata
        """
        if not file_path and not file_content:
            return ExtractionResult(
                success=False,
                text='',
                metadata={},
                error='Either file_path or file_content must be provided'
            )
        
        # Determine filename
        if not filename:
            if file_path:
                filename = Path(file_path).name
            else:
                filename = 'unknown'
        
        ext = self._get_extension(filename)
        
        if not ext:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'extension': None},
                error='Cannot determine file type (no extension)'
            )
        
        if ext not in self.SUPPORTED_EXTENSIONS:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'extension': ext},
                error=f'File type .{ext} is not supported'
            )
        
        # Read file content if needed
        if file_content is None and file_path:
            try:
                with open(file_path, 'rb') as f:
                    file_content = f.read()
            except Exception as e:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata={'filename': filename, 'extension': ext},
                    error=f'Failed to read file: {str(e)}'
                )
        
        # Extract based on file type
        # Use subprocess for libraries in dedicated env, direct import for text files
        try:
            if ext == 'pdf':
                return self._extract_pdf(file_content, filename)
            elif ext in ['docx', 'doc']:
                return self._extract_docx(file_content, filename)
            elif ext in ['xlsx', 'xls']:
                return self._extract_excel(file_content, filename, ext)
            elif ext == 'pptx':
                return self._extract_pptx(file_content, filename)
            elif ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'webp']:
                return self._extract_image(file_content, filename, ext)
            elif ext in ['txt', 'md', 'py', 'js', 'ts', 'css', 'java', 'cpp', 'c', 'h', 'rs', 'go']:
                return self._extract_text(file_content, filename)
            elif ext == 'csv':
                return self._extract_csv(file_content, filename)
            elif ext == 'json':
                return self._extract_json(file_content, filename)
            elif ext in ['html', 'xml']:
                return self._extract_markup(file_content, filename, ext)
            elif ext in ['yaml', 'yml']:
                return self._extract_text(file_content, filename)
            else:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata={'filename': filename, 'extension': ext},
                    error=f'Extraction not implemented for .{ext}'
                )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'extension': ext},
                error=f'Extraction error: {str(e)}'
            )
    
    def _extract_pdf(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from PDF"""
        text_parts = []
        metadata = {'filename': filename, 'type': 'PDF', 'pages': 0}
        
        # Use subprocess if dedicated env is available
        python_exe = self._get_python_executable()
        if python_exe and self._libraries_available['pymupdf']:
            return self._extract_via_subprocess('pdf', content, filename, 'pdf', metadata)
        
        # Try PyMuPDF first (best quality) - direct import
        if self._libraries_available['pymupdf']:
            try:
                import fitz
                doc = fitz.open(stream=content, filetype='pdf')
                for page_num, page in enumerate(doc, 1):
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"--- Page {page_num} ---\n{text}")
                metadata['pages'] = len(doc)
                doc.close()
                return ExtractionResult(
                    success=True,
                    text='\n\n'.join(text_parts),
                    metadata=metadata
                )
            except Exception as e:
                pass  # Fall back to pypdf
        
        # Try pypdf/PyPDF2
        if self._libraries_available['pypdf']:
            try:
                try:
                    import pypdf
                    reader = pypdf.PdfReader(io.BytesIO(content))
                except ImportError:
                    import PyPDF2
                    reader = PyPDF2.PdfReader(io.BytesIO(content))
                
                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(f"--- Page {page_num} ---\n{text}")
                metadata['pages'] = len(reader.pages)
                
                return ExtractionResult(
                    success=True,
                    text='\n\n'.join(text_parts),
                    metadata=metadata
                )
            except Exception as e:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata=metadata,
                    error=f'PDF extraction failed: {str(e)}. Install PyMuPDF (pip install pymupdf) or pypdf (pip install pypdf)'
                )
        
        return ExtractionResult(
            success=False,
            text='',
            metadata=metadata,
            error='No PDF library available. Setup document extraction environment first.'
        )
    
    def _extract_via_subprocess(self, file_type: str, content: bytes, filename: str, ext: str, metadata: Dict) -> ExtractionResult:
        """Extract using subprocess with dedicated environment"""
        import base64
        import tempfile
        import os
        
        def get_short_path(path):
            """Get Windows short path to avoid path length issues"""
            if os.name == 'nt':  # Windows
                try:
                    import ctypes
                    from ctypes import wintypes
                    GetShortPathNameW = ctypes.windll.kernel32.GetShortPathNameW
                    GetShortPathNameW.argtypes = [wintypes.LPCWSTR, wintypes.LPWSTR, wintypes.DWORD]
                    GetShortPathNameW.restype = wintypes.DWORD
                    
                    buf = ctypes.create_unicode_buffer(260)
                    result = GetShortPathNameW(str(path), buf, 260)
                    if result:
                        return Path(buf.value)
                except Exception:
                    pass
            return path
        
        try:
            # Create temporary script file (more reliable than inline script)
            script_content = '''import sys
import base64
import io
import json

content_b64 = sys.argv[1]
file_type = sys.argv[2]
filename = sys.argv[3]

try:
    content = base64.b64decode(content_b64)
except Exception as e:
    print(json.dumps({'success': False, 'error': f'Failed to decode content: {str(e)}'}))
    sys.exit(1)

try:
    if file_type == 'pdf':
        import fitz
        doc = fitz.open(stream=content, filetype='pdf')
        text_parts = []
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"--- Page {page_num} ---\\n{text}")
        result = {
            'success': True,
            'text': '\\n\\n'.join(text_parts),
            'metadata': {'pages': len(doc)}
        }
        doc.close()
    elif file_type == 'docx':
        from docx import Document
        doc = Document(io.BytesIO(content))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                text_parts.append('\\n--- Table ---\\n' + '\\n'.join(table_text))
        result = {
            'success': True,
            'text': '\\n\\n'.join(text_parts),
            'metadata': {'paragraphs': len(doc.paragraphs)}
        }
    elif file_type == 'xlsx':
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    sheet_text.append(row_text)
            if len(sheet_text) > 1:
                text_parts.append('\\n'.join(sheet_text))
        result = {
            'success': True,
            'text': '\\n\\n'.join(text_parts),
            'metadata': {'sheets': wb.sheetnames}
        }
    else:
        result = {'success': False, 'error': f'Type {file_type} not implemented in subprocess'}
except Exception as e:
    result = {'success': False, 'error': str(e)}

print(json.dumps(result))
'''
            
            # Use shorter temp directory to avoid Windows path length issues
            temp_dir = Path(tempfile.gettempdir()) / 'beep_ocr'
            temp_dir.mkdir(exist_ok=True)
            
            # Create a short filename
            import hashlib
            file_hash = hashlib.md5(content[:100]).hexdigest()[:8]
            script_path = temp_dir / f'ocr_{file_hash}.py'
            
            # Write script to temp file
            with open(script_path, 'w', encoding='utf-8') as f:
                f.write(script_content)
            
            # Get short path on Windows to avoid path length issues
            script_path_short = get_short_path(script_path)
            python_exe = self._get_python_executable()
            if not python_exe:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata=metadata,
                    error='Python executable not found. Please set up the document extraction environment.'
                )
            python_exe_short = get_short_path(python_exe)
            
            try:
                # Validate Python executable exists
                if not python_exe_short.exists():
                    return ExtractionResult(
                        success=False,
                        text='',
                        metadata=metadata,
                        error=f'Python executable not found: {python_exe_short}. Please set up the document extraction environment.'
                    )
                
                # Run in dedicated environment using short paths
                content_b64 = base64.b64encode(content).decode('utf-8')
                logger.info(f"Running subprocess: {python_exe_short} {script_path_short} [content] {file_type} {filename}")
                
                result = subprocess.run(
                    [str(python_exe_short), str(script_path_short), content_b64, file_type, filename],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                
                if result.returncode == 0:
                    import json
                    data = json.loads(result.stdout)
                    metadata.update(data.get('metadata', {}))
                    return ExtractionResult(
                        success=data.get('success', False),
                        text=data.get('text', ''),
                        metadata=metadata,
                        error=data.get('error')
                    )
                else:
                    error_msg = result.stderr[:500] if result.stderr else 'Unknown error'
                    logger.error(f"Subprocess extraction failed for {filename}: {error_msg}")
                    if result.stdout:
                        logger.debug(f"Subprocess stdout: {result.stdout[:500]}")
                    return ExtractionResult(
                        success=False,
                        text='',
                        metadata=metadata,
                        error=f'Subprocess extraction failed: {error_msg}'
                    )
            finally:
                # Clean up temp file
                try:
                    os.unlink(script_path)
                except:
                    pass
        except subprocess.TimeoutExpired:
            logger.error(f"Subprocess timeout for {filename} after 60 seconds")
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error='Extraction timed out after 60 seconds. The file may be too large or complex.'
            )
        except Exception as e:
            logger.error(f"Subprocess exception for {filename}: {e}", exc_info=True)
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error=f'Subprocess error: {str(e)}'
            )
    
    def _extract_docx(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from DOCX"""
        metadata = {'filename': filename, 'type': 'Word'}
        
        # Use subprocess if dedicated env is available
        python_exe = self._get_python_executable()
        if python_exe and self._libraries_available['python_docx']:
            return self._extract_via_subprocess('docx', content, filename, ext, metadata)
        
        if not self._libraries_available['python_docx']:
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error='python-docx not installed. Setup document extraction environment first.'
            )
        
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_text.append(row_text)
                if table_text:
                    text_parts.append('\n--- Table ---\n' + '\n'.join(table_text))
            
            return ExtractionResult(
                success=True,
                text='\n\n'.join(text_parts),
                metadata={'filename': filename, 'type': 'Word', 'paragraphs': len(doc.paragraphs)}
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error=f'DOCX extraction failed: {str(e)}'
            )
    
    def _extract_excel(self, content: bytes, filename: str, ext: str) -> ExtractionResult:
        """Extract text from Excel files"""
        text_parts = []
        metadata = {'filename': filename, 'type': 'Excel', 'sheets': []}
        
        # Use subprocess if dedicated env is available
        python_exe = self._get_python_executable()
        if python_exe and self._libraries_available['openpyxl']:
            return self._extract_via_subprocess('xlsx', content, filename, ext, metadata)
        
        # Try openpyxl for .xlsx
        if ext == 'xlsx' and self._libraries_available['openpyxl']:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(content), data_only=True)
                
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_text = [f"--- Sheet: {sheet_name} ---"]
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                        if row_text.strip():
                            sheet_text.append(row_text)
                    
                    if len(sheet_text) > 1:
                        text_parts.append('\n'.join(sheet_text))
                        metadata['sheets'].append(sheet_name)
                
                return ExtractionResult(
                    success=True,
                    text='\n\n'.join(text_parts),
                    metadata=metadata
                )
            except Exception as e:
                pass  # Fall back to xlrd
        
        # Try xlrd for .xls or as fallback
        if self._libraries_available['xlrd']:
            try:
                import xlrd
                book = xlrd.open_workbook(file_contents=content)
                
                for sheet_idx, sheet in enumerate(book.sheets()):
                    sheet_text = [f"--- Sheet: {sheet.name} ---"]
                    
                    for row_idx in range(sheet.nrows):
                        row = sheet.row(row_idx)
                        row_text = ' | '.join(str(cell.value) for cell in row)
                        if row_text.strip():
                            sheet_text.append(row_text)
                    
                    if len(sheet_text) > 1:
                        text_parts.append('\n'.join(sheet_text))
                        metadata['sheets'].append(sheet.name)
                
                return ExtractionResult(
                    success=True,
                    text='\n\n'.join(text_parts),
                    metadata=metadata
                )
            except Exception as e:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata=metadata,
                    error=f'Excel extraction failed: {str(e)}'
                )
        
        return ExtractionResult(
            success=False,
            text='',
            metadata=metadata,
            error='No Excel library available. Setup document extraction environment first.'
        )
    
    def _extract_image(self, content: bytes, filename: str, ext: str) -> ExtractionResult:
        """Extract text from images using OCR"""
        metadata = {'filename': filename, 'type': 'Image', 'format': ext.upper()}
        
        # Refresh library check before attempting extraction (in case packages were just installed)
        python_exe = self._get_python_executable()
        if python_exe:
            # Re-check easyocr availability in the dedicated environment
            try:
                result = subprocess.run(
                    [python_exe, '-c', 'import easyocr; print("ok")'],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                easyocr_available = result.returncode == 0
                if easyocr_available:
                    self._libraries_available['easyocr'] = True
                    logger.info(f"EasyOCR confirmed available in dedicated environment: {python_exe}")
                else:
                    logger.warning(f"EasyOCR not available in dedicated environment: {result.stderr[:200]}")
                    self._libraries_available['easyocr'] = False
            except Exception as e:
                logger.warning(f"Failed to check EasyOCR in dedicated environment: {e}")
                easyocr_available = self._libraries_available.get('easyocr', False)
        else:
            easyocr_available = self._libraries_available.get('easyocr', False)
        
        # Use subprocess if dedicated env is available
        if python_exe and easyocr_available:
            logger.info(f"Using dedicated environment for OCR: {python_exe}")
            return self._extract_image_via_subprocess(content, filename, ext, metadata)
        
        if not easyocr_available:
            error_msg = 'OCR not available. Install easyocr package. Setup document extraction environment first.'
            python_exe = self._get_python_executable()
            if python_exe:
                error_msg += f' (Checked in: {python_exe})'
            logger.error(f"OCR extraction failed for {filename}: {error_msg}")
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error=error_msg
            )
        
        try:
            import easyocr
            from PIL import Image
            import numpy as np
            
            # Try to import OpenCV for image preprocessing (optional but improves quality)
            try:
                import cv2
                has_opencv = True
            except ImportError:
                has_opencv = False
            
            # Load image from bytes
            image = Image.open(io.BytesIO(content))
            
            # Get image metadata
            metadata['width'] = image.width
            metadata['height'] = image.height
            metadata['mode'] = image.mode
            
            # Preprocess image to improve OCR accuracy
            if has_opencv:
                # Convert to RGB if needed
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Convert to numpy array for OpenCV processing
                img_array = np.array(image)
                
                # Convert RGB to BGR for OpenCV
                img_bgr = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
                
                # Convert to grayscale for preprocessing
                gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)
                
                # Apply image enhancement techniques
                # 1. Denoise
                denoised = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
                
                # 2. Increase contrast using CLAHE
                clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
                enhanced = clahe.apply(denoised)
                
                # 3. Apply sharpening filter
                kernel = np.array([[-1, -1, -1],
                                   [-1,  9, -1],
                                   [-1, -1, -1]])
                sharpened = cv2.filter2D(enhanced, -1, kernel)
                
                # 4. Threshold to get binary image
                _, thresh = cv2.threshold(sharpened, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
                
                # Convert back to RGB for EasyOCR
                processed_img = cv2.cvtColor(thresh, cv2.COLOR_GRAY2RGB)
            else:
                # Fallback: basic preprocessing without OpenCV
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                processed_img = np.array(image)
            
            # Initialize EasyOCR reader (lazy load, cache it)
            if not hasattr(self, '_ocr_reader'):
                # Use English by default, can be extended to support more languages
                self._ocr_reader = easyocr.Reader(['en'], gpu=False, verbose=False)
            
            # Perform OCR with better parameters
            results = self._ocr_reader.readtext(processed_img, paragraph=False, detail=1)
            
            # Combine all detected text
            text_parts = [result[1] for result in results]  # result[1] is the text
            text = '\n'.join(text_parts)
            
            # Add OCR confidence scores to metadata
            if results:
                avg_confidence = sum(result[2] for result in results) / len(results)
                metadata['ocr_confidence'] = round(avg_confidence * 100, 2)
                metadata['text_regions'] = len(results)
            
            return ExtractionResult(
                success=True,
                text=text.strip(),
                metadata=metadata
            )
        except ImportError as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error=f'OCR library not available: {str(e)}. Install easyocr and Pillow.'
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error=f'OCR extraction failed: {str(e)}'
            )
    
    def _extract_image_via_subprocess(self, content: bytes, filename: str, ext: str, metadata: Dict, ocr_engine=None) -> ExtractionResult:
        """Extract text from images using subprocess with dedicated environment"""
        import base64
        import tempfile
        import os
        
        def get_short_path(path):
            """Get Windows short path to avoid path length issues"""
            if os.name == 'nt':  # Windows
                try:
                    import ctypes
                    from ctypes import wintypes
                    GetShortPathNameW = ctypes.windll.kernel32.GetShortPathNameW
                    GetShortPathNameW.argtypes = [wintypes.LPCWSTR, wintypes.LPWSTR, wintypes.DWORD]
                    GetShortPathNameW.restype = wintypes.DWORD
                    
                    buf = ctypes.create_unicode_buffer(260)
                    result = GetShortPathNameW(str(path), buf, 260)
                    if result:
                        return Path(buf.value)
                except Exception:
                    pass
            return path
        
        def get_short_temp_dir():
            """Get a short temporary directory path to avoid Windows path length issues"""
            if os.name == 'nt':  # Windows
                # Use C:\temp or C:\TEMP for shortest path
                short_temp = Path('C:/temp')
                if not short_temp.exists():
                    try:
                        short_temp.mkdir(parents=True, exist_ok=True)
                    except:
                        # Fallback to user temp
                        short_temp = Path(tempfile.gettempdir())
                return short_temp
            return Path(tempfile.gettempdir())
        
        try:
            script_content = '''import sys
import json
from pathlib import Path
import warnings
import os

# Suppress warnings to avoid polluting stdout
warnings.filterwarnings('ignore')
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'  # Suppress TensorFlow warnings

# Redirect stderr to avoid warnings in stdout
import sys
sys.stderr = open(os.devnull, 'w')

# Get arguments - now receives image file path instead of base64 (avoids command line length issues)
image_file_path = sys.argv[1]
filename = sys.argv[2]
ext = sys.argv[3]

# Try to import OpenCV for image preprocessing (optional but improves quality)
try:
    import cv2
    has_opencv = True
except ImportError:
    has_opencv = False

# Load image from file path (much shorter than base64 in command line)
from PIL import Image
import numpy as np
image = Image.open(image_file_path)

# Get image metadata
width = image.width
height = image.height
mode = image.mode

# Preprocess image to improve OCR accuracy
if has_opencv:
    # Convert to RGB if needed
    if image.mode != 'RGB':
        image = image.convert('RGB')
    
    # Convert to numpy array for OpenCV processing
    img_array = np.array(image)
    
    # Convert RGB to BGR for OpenCV
    img_bgr = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
    
    # Convert to grayscale for preprocessing
    gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)
    
    # Apply image enhancement techniques
    # 1. Denoise
    denoised = cv2.fastNlMeansDenoising(gray, None, 10, 7, 21)
    
    # 2. Increase contrast using CLAHE
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    enhanced = clahe.apply(denoised)
    
    # 3. Apply sharpening filter
    kernel = np.array([[-1, -1, -1],
                       [-1,  9, -1],
                       [-1, -1, -1]])
    sharpened = cv2.filter2D(enhanced, -1, kernel)
    
    # 4. Threshold to get binary image
    _, thresh = cv2.threshold(sharpened, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    
    # Prepare processed image based on OCR engine
    if engine_type == 'tesseract':
        processed_img = Image.fromarray(thresh)
    elif engine_type == 'paddleocr':
        processed_img = cv2.cvtColor(thresh, cv2.COLOR_GRAY2BGR)
    else:  # easyocr
        processed_img = cv2.cvtColor(thresh, cv2.COLOR_GRAY2RGB)
else:
    # Fallback: basic preprocessing without OpenCV
    if image.mode != 'RGB':
        image = image.convert('RGB')
    processed_img = np.array(image) if engine_type != 'tesseract' else image.convert('L')

# Perform OCR based on selected engine
try:
    if engine_type == 'easyocr':
        import easyocr
        reader = easyocr.Reader(['en'], gpu=False, verbose=False)
        results = reader.readtext(processed_img, paragraph=False, detail=1)
        text_parts = [result[1] for result in results]
        text = '\\n'.join(text_parts)
        if results:
            avg_confidence = sum(result[2] for result in results) / len(results)
            metadata = {'width': width, 'height': height, 'mode': mode, 'ocr_confidence': round(avg_confidence * 100, 2), 'text_regions': len(results)}
        else:
            metadata = {'width': width, 'height': height, 'mode': mode}
    
    elif engine_type == 'tesseract':
        import pytesseract
        text = pytesseract.image_to_string(processed_img, lang='eng', config='--psm 6')
        try:
            data = pytesseract.image_to_data(processed_img, lang='eng', output_type=pytesseract.Output.DICT)
            confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
            if confidences:
                avg_confidence = sum(confidences) / len(confidences)
                metadata = {'width': width, 'height': height, 'mode': mode, 'ocr_confidence': round(avg_confidence, 2), 'text_regions': len(confidences)}
            else:
                metadata = {'width': width, 'height': height, 'mode': mode}
        except:
            metadata = {'width': width, 'height': height, 'mode': mode}
    
    elif engine_type == 'paddleocr':
        from paddleocr import PaddleOCR
        reader = PaddleOCR(use_angle_cls=True, lang='en', use_gpu=False)
        results = reader.ocr(processed_img, cls=True)
        text_parts = []
        confidences = []
        if results and results[0]:
            for line in results[0]:
                if line and len(line) >= 2:
                    text_info = line[1]
                    if isinstance(text_info, tuple) and len(text_info) >= 2:
                        text_parts.append(text_info[0])
                        confidences.append(text_info[1])
        text = '\\n'.join(text_parts)
        if confidences:
            avg_confidence = sum(confidences) / len(confidences)
            metadata = {'width': width, 'height': height, 'mode': mode, 'ocr_confidence': round(avg_confidence * 100, 2), 'text_regions': len(confidences)}
        else:
            metadata = {'width': width, 'height': height, 'mode': mode}
    
    else:
        raise ValueError(f'Unknown OCR engine: {engine_type}')
    
    result = {
        'success': True,
        'text': text.strip(),
        'metadata': metadata
    }
except Exception as e:
    result = {'success': False, 'error': str(e)}

# Output only JSON to stdout (single line, no extra output)
print(json.dumps(result), flush=True)
'''
            
            # Use short temp directory to avoid Windows path length issues
            temp_dir = get_short_temp_dir() / 'beep_ocr'
            temp_dir.mkdir(exist_ok=True)
            
            # Create very short filenames to minimize path length
            import hashlib
            import time
            timestamp = int(time.time() * 1000) % 100000  # Last 5 digits
            script_path = temp_dir / f'o{timestamp}.py'  # Very short: o12345.py
            
            # Write image to temp file instead of passing base64 via command line
            # This avoids Windows command line length limits
            image_hash = hashlib.md5(content).hexdigest()[:8]
            image_path = temp_dir / f'i{timestamp}.{ext}'  # Very short: i12345.png
            
            try:
                # Write script to temp file
                with open(script_path, 'w', encoding='utf-8') as f:
                    f.write(script_content)
                
                # Write image to temp file
                with open(image_path, 'wb') as f:
                    f.write(content)
            except Exception as e:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata=metadata,
                    error=f'Failed to create temporary files: {str(e)}'
                )
            
            # Get short paths on Windows to avoid path length issues
            try:
                script_path_short = get_short_path(script_path)
            except:
                script_path_short = script_path
            
            try:
                image_path_short = get_short_path(image_path)
            except:
                image_path_short = image_path
            
            python_exe = self._get_python_executable()
            if not python_exe:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata=metadata,
                    error='Python executable not found. Please set up the document extraction environment.'
                )
            
            try:
                python_exe_short = get_short_path(python_exe)
            except:
                python_exe_short = python_exe
            
            try:
                # Validate Python executable exists
                if not python_exe_short.exists():
                    return ExtractionResult(
                        success=False,
                        text='',
                        metadata=metadata,
                        error=f'Python executable not found: {python_exe_short}. Please set up the document extraction environment.'
                    )
                
                # Run in dedicated environment using short paths
                # Pass image file path instead of base64 content (avoids command line length limits)
                logger.info(f"Running OCR subprocess ({engine_type}): {python_exe_short} {script_path_short} {image_path_short} {filename} {ext} {engine_type}")
                
                result = subprocess.run(
                    [str(python_exe_short), str(script_path_short), str(image_path_short), filename, ext, engine_type],
                    capture_output=True,
                    text=True,
                    timeout=120,  # OCR can take longer
                    cwd=str(temp_dir)  # Set working directory to minimize path lengths
                )
                
                if result.returncode == 0:
                    import json
                    # Extract JSON from stdout (handle cases where there might be extra output)
                    stdout_lines = result.stdout.strip().split('\n')
                    json_line = None
                    
                    # Find the line that contains JSON (usually the last line)
                    for line in reversed(stdout_lines):
                        line = line.strip()
                        if line and (line.startswith('{') or line.startswith('[')):
                            json_line = line
                            break
                    
                    if not json_line:
                        # If no JSON found, try parsing the entire stdout
                        json_line = result.stdout.strip()
                    
                    try:
                        data = json.loads(json_line)
                        metadata.update(data.get('metadata', {}))
                        return ExtractionResult(
                            success=data.get('success', False),
                            text=data.get('text', ''),
                            metadata=metadata,
                            error=data.get('error')
                        )
                    except json.JSONDecodeError as e:
                        # Log the actual output for debugging
                        logger.error(f"Failed to parse OCR subprocess JSON output. stdout: {result.stdout[:500]}, stderr: {result.stderr[:500]}")
                        return ExtractionResult(
                            success=False,
                            text='',
                            metadata=metadata,
                            error=f'OCR subprocess returned invalid JSON: {str(e)}. Output: {result.stdout[:200]}'
                        )
                else:
                    error_msg = result.stderr[:500] if result.stderr else result.stdout[:500] or 'Unknown error'
                    logger.error(f"OCR subprocess failed with return code {result.returncode}. stderr: {error_msg}")
                    return ExtractionResult(
                        success=False,
                        text='',
                        metadata=metadata,
                        error=f'OCR subprocess failed (exit code {result.returncode}): {error_msg}'
                    )
            finally:
                # Clean up temp files
                try:
                    if script_path.exists():
                        os.unlink(script_path)
                except Exception:
                    pass
                try:
                    if 'image_path' in locals() and image_path.exists():
                        os.unlink(image_path)
                except Exception:
                    pass
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata=metadata,
                error=f'OCR subprocess error: {str(e)}'
            )
    
    def _extract_pptx(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from PowerPoint"""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if len(slide_text) > 1:
                    text_parts.append('\n'.join(slide_text))
            
            return ExtractionResult(
                success=True,
                text='\n\n'.join(text_parts),
                metadata={'filename': filename, 'type': 'PowerPoint', 'slides': len(prs.slides)}
            )
        except ImportError:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'type': 'PowerPoint'},
                error='python-pptx not installed. Install with: pip install python-pptx'
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'type': 'PowerPoint'},
                error=f'PPTX extraction failed: {str(e)}'
            )
    
    def _extract_text(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from plain text files"""
        try:
            # Try UTF-8 first
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                # Fall back to latin-1
                text = content.decode('latin-1')
            except UnicodeDecodeError:
                return ExtractionResult(
                    success=False,
                    text='',
                    metadata={'filename': filename},
                    error='Failed to decode text file (not UTF-8 or latin-1)'
                )
        
        return ExtractionResult(
            success=True,
            text=text,
            metadata={'filename': filename, 'type': 'Text', 'chars': len(text)}
        )
    
    def _extract_csv(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from CSV"""
        try:
            import csv
            text = content.decode('utf-8')
            reader = csv.reader(io.StringIO(text))
            
            rows = []
            for i, row in enumerate(reader, 1):
                if i > 1000:  # Limit rows
                    rows.append("... (truncated)")
                    break
                rows.append(' | '.join(row))
            
            return ExtractionResult(
                success=True,
                text='\n'.join(rows),
                metadata={'filename': filename, 'type': 'CSV', 'rows': min(len(rows), 1000)}
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'type': 'CSV'},
                error=f'CSV extraction failed: {str(e)}'
            )
    
    def _extract_json(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from JSON"""
        try:
            import json
            text = content.decode('utf-8')
            data = json.loads(text)
            
            # Format JSON nicely
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            
            return ExtractionResult(
                success=True,
                text=formatted,
                metadata={'filename': filename, 'type': 'JSON'}
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'type': 'JSON'},
                error=f'JSON extraction failed: {str(e)}'
            )
    
    def _extract_markup(self, content: bytes, filename: str, ext: str) -> ExtractionResult:
        """Extract text from HTML/XML"""
        try:
            from html.parser import HTMLParser
            
            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text = []
                    self.ignore_tags = {'script', 'style', 'meta', 'link'}
                
                def handle_data(self, data):
                    if data.strip():
                        self.text.append(data.strip())
                
                def handle_starttag(self, tag, attrs):
                    if tag in self.ignore_tags:
                        self.ignore = True
                
                def handle_endtag(self, tag):
                    self.ignore = False
            
            text = content.decode('utf-8')
            parser = TextExtractor()
            parser.feed(text)
            
            return ExtractionResult(
                success=True,
                text='\n'.join(parser.text),
                metadata={'filename': filename, 'type': ext.upper()}
            )
        except Exception as e:
            return ExtractionResult(
                success=False,
                text='',
                metadata={'filename': filename, 'type': ext.upper()},
                error=f'{ext.upper()} extraction failed: {str(e)}'
            )
    
    def get_required_packages(self) -> List[str]:
        """Get list of recommended packages for full extraction support"""
        packages = []
        
        if not self._libraries_available['pypdf'] and not self._libraries_available['pymupdf']:
            packages.append('pymupdf')  # Prefer PyMuPDF for better quality
        
        if not self._libraries_available['python_docx']:
            packages.append('python-docx')
        
        if not self._libraries_available['openpyxl']:
            packages.append('openpyxl')
        
        if not self._libraries_available['xlrd']:
            packages.append('xlrd')
        
        # Optional but recommended
        packages.append('python-pptx')  # For PowerPoint
        
        # OCR packages
        # OCR packages
        if not self._libraries_available['easyocr']:
            packages.append('easyocr  # For OCR - Pure Python, no external dependencies')
        
        if not self._libraries_available['tesseract']:
            packages.append('pytesseract  # For Tesseract OCR - Highly accurate for printed text')
        
        if not self._libraries_available['paddleocr']:
            packages.append('paddleocr  # For PaddleOCR - Fast and accurate, excellent for Chinese/English')
        
        if not self._libraries_available['pillow']:
            packages.append('Pillow  # For image processing')
        
        # OpenCV for image preprocessing (improves OCR quality)
        packages.append('opencv-python  # For image preprocessing to improve OCR accuracy')
        
        return packages
    
    def get_status(self) -> Dict[str, Any]:
        """Get status of extraction libraries"""
        # Group formats by category
        formats_by_category = {
            'Documents': ['pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt'],
            'Images (OCR)': ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'webp'],
            'Text': ['txt', 'md', 'csv', 'json', 'html', 'xml', 'yaml', 'yml'],
            'Code': ['py', 'js', 'ts', 'css', 'java', 'cpp', 'c', 'h', 'rs', 'go']
        }
        
        return {
            'libraries': self._libraries_available,
            'supported_formats': list(self.SUPPORTED_EXTENSIONS.keys()),
            'formats_by_category': formats_by_category,
            'recommended_packages': self.get_required_packages(),
            'using_dedicated_env': self.use_dedicated_env and self.env_mgr is not None,
            'env_python': str(self._get_python_executable()) if self._get_python_executable() else None,
            'ocr_available': self._libraries_available.get('easyocr', False) and self._libraries_available.get('pillow', False)
        }


def get_document_extractor() -> DocumentExtractor:
    """Get singleton instance of DocumentExtractor"""
    if not hasattr(get_document_extractor, '_instance'):
        get_document_extractor._instance = DocumentExtractor()
    return get_document_extractor._instance


def get_document_extractor() -> DocumentExtractor:
    """Get singleton instance of DocumentExtractor"""
    if not hasattr(get_document_extractor, '_instance'):
        get_document_extractor._instance = DocumentExtractor()
    return get_document_extractor._instance
